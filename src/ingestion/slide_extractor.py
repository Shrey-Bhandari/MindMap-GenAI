from pptx import Presentation
from src.ingestion.normalize_text import normalize_text

def extract_slide_text(ppt_path: str) -> list[dict]:
    prs = Presentation(ppt_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    texts.append(t)

        combined = "\n".join(texts)
        combined = normalize_text(combined)

        if combined:  # Skip empty slides
            slides_data.append({
                "slide_id": idx,
                "text": combined
            })

    return slides_data
